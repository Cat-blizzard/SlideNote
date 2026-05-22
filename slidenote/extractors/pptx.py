from __future__ import annotations

from pathlib import Path
from typing import Iterable

from slidenote.image_assets import image_metadata, refine_image_role_for_placement
from slidenote.models import Deck, ImageAsset, SlidePage, TableBlock, TextBlock, normalize_rel_path
from slidenote.rendering import render_pptx_screenshots
from slidenote.utils import unique_path


def extract_pptx(input_path: Path, output_root: Path) -> Deck:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PPTX extraction. Install with `pip install python-pptx`.") from exc

    images_dir = output_root / "images"
    screenshots_dir = output_root / "screenshots"
    images_dir.mkdir(parents=True, exist_ok=True)
    screenshots_dir.mkdir(parents=True, exist_ok=True)

    presentation = Presentation(str(input_path))
    screenshot_map, render_warnings = render_pptx_screenshots(input_path, screenshots_dir, output_root)
    pages: list[SlidePage] = []
    slide_width = float(presentation.slide_width)
    slide_height = float(presentation.slide_height)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        page = SlidePage(slide_id=slide_index, page_width=slide_width, page_height=slide_height)
        counters = {"text": 1, "table": 1, "image": 1}
        for shape in _iter_shapes(slide.shapes):
            if getattr(shape, "has_table", False):
                table = _extract_table(shape, slide_index, counters["table"])
                if table:
                    page.tables.append(table)
                    counters["table"] += 1
                continue

            if _is_picture(shape):
                image = _extract_picture(
                    shape,
                    slide_index,
                    counters["image"],
                    images_dir,
                    output_root,
                    slide_width=slide_width,
                    slide_height=slide_height,
                )
                if image:
                    page.images.append(image)
                    counters["image"] += 1
                continue

            text, style_runs = _shape_text_and_runs(shape)
            if text:
                block_type = _classify_shape_text(shape, text)
                block = TextBlock(
                    id=f"s{slide_index}_t{counters['text']}",
                    type=block_type,
                    content=text,
                    bbox=_shape_bbox(shape),
                    style_runs=style_runs,
                )
                page.text_blocks.append(block)
                counters["text"] += 1
                if not page.title and block_type == "title":
                    page.title = text.splitlines()[0].strip()

        if not page.title:
            page.title = _fallback_title(page.text_blocks)
        page.page_screenshot = screenshot_map.get(slide_index)
        page.notes = _notes_text(slide)
        if page.page_screenshot is None:
            page.warnings.append("No full-slide screenshot available. Install LibreOffice or PowerPoint for screenshot rendering.")
        pages.append(page)

    return Deck(source_path=str(input_path), source_type="pptx", pages=pages, warnings=render_warnings)


def _iter_shapes(shapes: Iterable[object]) -> Iterable[object]:
    for shape in shapes:
        if hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _shape_text(shape: object) -> str:
    text, _ = _shape_text_and_runs(shape)
    return text


def _shape_text_and_runs(shape: object) -> tuple[str, list[dict[str, object]]]:
    if not getattr(shape, "has_text_frame", False):
        return "", []
    frame = shape.text_frame
    lines: list[str] = []
    style_runs: list[dict[str, object]] = []
    has_explicit_style = False
    for paragraph in frame.paragraphs:
        runs = [(run.text, _run_style(run)) for run in paragraph.runs if run.text]
        text = "".join(run_text for run_text, _ in runs).strip()
        if text:
            indent = "  " * int(getattr(paragraph, "level", 0) or 0)
            if lines:
                style_runs.append({"text": "\n"})
            if indent:
                style_runs.append({"text": indent})
            lines.append(f"{indent}{text}")
            for run_text, style in runs:
                if not run_text:
                    continue
                entry: dict[str, object] = {"text": run_text}
                if style:
                    entry.update(style)
                    has_explicit_style = True
                style_runs.append(entry)
    return "\n".join(lines).strip(), style_runs if has_explicit_style else []


def _run_style(run: object) -> dict[str, object]:
    style: dict[str, object] = {}
    color = _run_font_color(run)
    if color:
        style["color"] = color
    try:
        bold = run.font.bold
    except Exception:
        bold = None
    if bold is True:
        style["bold"] = True
    try:
        italic = run.font.italic
    except Exception:
        italic = None
    if italic is True:
        style["italic"] = True
    return style


def _run_font_color(run: object) -> str | None:
    try:
        rgb = run.font.color.rgb
    except Exception:
        return None
    if not rgb:
        return None
    value = str(rgb).strip()
    if not value:
        return None
    if len(value) == 6 and all(char in "0123456789abcdefABCDEF" for char in value):
        return f"#{value.upper()}"
    return None


def _classify_shape_text(shape: object, text: str) -> str:
    placeholder_type = None
    if getattr(shape, "is_placeholder", False):
        try:
            placeholder_type = str(shape.placeholder_format.type).upper()
        except Exception:
            placeholder_type = None
    if placeholder_type and "TITLE" in placeholder_type:
        return "title"
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if len(lines) > 1:
        return "bullet"
    if len(text) <= 120:
        return "heading"
    return "paragraph"


def _extract_table(shape: object, slide_index: int, table_index: int) -> TableBlock | None:
    rows: list[list[str]] = []
    for row in shape.table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    if not any(any(cell for cell in row) for row in rows):
        return None
    return TableBlock(id=f"s{slide_index}_tbl{table_index}", rows=rows, bbox=_shape_bbox(shape))


def _is_picture(shape: object) -> bool:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return hasattr(shape, "image")


def _extract_picture(
    shape: object,
    slide_index: int,
    image_index: int,
    images_dir: Path,
    output_root: Path,
    slide_width: float,
    slide_height: float,
) -> ImageAsset | None:
    if not hasattr(shape, "image"):
        return None
    image = shape.image
    ext = image.ext or "png"
    image_path = unique_path(images_dir / f"slide{slide_index}_img{image_index}.{ext}")
    image_path.write_bytes(image.blob)
    meta = image_metadata(image_path)
    bbox = _shape_bbox(shape)
    page_like = _is_page_like_shape(bbox, slide_width, slide_height)
    role = "page_image" if page_like else meta["role"]
    ignored = True if page_like else meta["ignored"]
    ignore_reason = "full_page_image" if page_like else meta["ignore_reason"]
    role, ignored, ignore_reason = refine_image_role_for_placement(
        role,
        ignored,
        ignore_reason,
        _shape_area_ratio(bbox, slide_width, slide_height),
        _shape_near_slide_edge(bbox, slide_width, slide_height),
    )
    return ImageAsset(
        id=f"s{slide_index}_img{image_index}",
        path=normalize_rel_path(image_path, output_root),
        caption=f"第 {slide_index} 页嵌入图片 {image_index}",
        bbox=bbox,
        source_format=ext,
        width=meta["width"],
        height=meta["height"],
        file_size=meta["file_size"],
        role=role,
        ignored=ignored,
        ignore_reason=ignore_reason,
    )


def _shape_bbox(shape: object) -> list[float] | None:
    try:
        return [float(shape.left), float(shape.top), float(shape.width), float(shape.height)]
    except Exception:
        return None


def _is_page_like_shape(bbox: list[float] | None, slide_width: float, slide_height: float) -> bool:
    if not bbox or slide_width <= 0 or slide_height <= 0:
        return False
    _, _, width, height = bbox
    return max(0.0, width) * max(0.0, height) / (slide_width * slide_height) >= 0.85


def _shape_area_ratio(bbox: list[float] | None, slide_width: float, slide_height: float) -> float | None:
    if not bbox or slide_width <= 0 or slide_height <= 0:
        return None
    _, _, width, height = bbox
    return max(0.0, width) * max(0.0, height) / (slide_width * slide_height)


def _shape_near_slide_edge(bbox: list[float] | None, slide_width: float, slide_height: float) -> bool:
    if not bbox or slide_width <= 0 or slide_height <= 0:
        return False
    left, top, width, height = bbox
    margin_x = slide_width * 0.08
    margin_y = slide_height * 0.08
    return left <= margin_x or top <= margin_y or left + width >= slide_width - margin_x or top + height >= slide_height - margin_y


def _fallback_title(blocks: list[TextBlock]) -> str | None:
    if not blocks:
        return None
    for block in blocks:
        if block.type in {"heading", "title"}:
            return block.content.splitlines()[0].strip()
    return blocks[0].content.splitlines()[0].strip()


def _notes_text(slide: object) -> str | None:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return None
    text = "\n".join(shape.text for shape in notes_slide.notes_text_frame.paragraphs if getattr(shape, "text", "").strip())
    return text.strip() or None
